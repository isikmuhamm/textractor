from typing import Optional, Set
from docx import Document
from odf import text, teletype
from odf.opendocument import load
import xlrd
from openpyxl import load_workbook
from pptx import Presentation
from .base_extractor import BaseExtractor
from ..utils.logger import logger

class OfficeExtractor(BaseExtractor):
    # Desteklenen office dosyası uzantıları
    SUPPORTED_EXTENSIONS: Set[str] = {
        '.docx', '.xlsx', '.pptx',
        '.odt', '.ods', '.odp',
        '.xls'
    }

    def can_handle(self, file_path: str) -> bool:
        """Dosya uzantısının desteklenip desteklenmediğini kontrol eder"""
        ext = file_path.lower().split('.')[-1] if '.' in file_path else ''
        return f'.{ext}' in self.SUPPORTED_EXTENSIONS

    def extract_text(self, file_path: str) -> Optional[str]:
        """Dosyadan text çıkarır"""
        ext = file_path.lower().split('.')[-1]
        
        try:
            if ext == 'docx':
                return self._extract_docx(file_path)
            elif ext == 'odt':
                return self._extract_odt(file_path)
            elif ext in ['xlsx', 'xls']:
                return self._extract_excel(file_path)
            elif ext == 'pptx':
                return self._extract_pptx(file_path)
            else:
                return None
        except Exception as e:
            logger.error(f"Office dosyası okuma hatası {file_path}: {str(e)}")
            return None

    def _extract_docx(self, file_path: str) -> str:
        """DOCX dosyasından text çıkarır"""
        doc = Document(file_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)

    def _extract_odt(self, file_path: str) -> str:
        """ODT dosyasından text çıkarır"""
        textdoc = load(file_path)
        allparas = textdoc.getElementsByType(text.P)
        return '\n'.join(teletype.extractText(para) for para in allparas)

    def _extract_excel(self, file_path: str) -> str:
        """Excel dosyasından text çıkarır"""
        if file_path.endswith('.xlsx'):
            wb = load_workbook(file_path, read_only=True)
            texts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                texts.append(f"\nSheet: {sheet}")
                for row in ws.rows:
                    row_texts = [str(cell.value) if cell.value is not None else '' for cell in row]
                    texts.append('\t'.join(row_texts))
            return '\n'.join(texts)
        else:  # .xls dosyaları için
            wb = xlrd.open_workbook(file_path)
            texts = []
            for sheet in wb.sheets():
                texts.append(f"\nSheet: {sheet.name}")
                for row in range(sheet.nrows):
                    row_texts = [str(cell.value) for cell in sheet.row(row)]
                    texts.append('\t'.join(row_texts))
            return '\n'.join(texts)

    def _extract_pptx(self, file_path: str) -> str:
        """PPTX dosyasından text çıkarır"""
        prs = Presentation(file_path)
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            texts.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return '\n'.join(texts)